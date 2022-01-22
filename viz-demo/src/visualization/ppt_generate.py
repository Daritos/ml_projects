from pptx import Presentation
from pptx.util import Inches
# https://python-pptx.readthedocs.io/


###
# Configure a new slide deck
###

prs=Presentation()

# Standard
prs.slide_width = Inches(4)
prs.slide_height = Inches(3)

''' 
# Widescreen
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
'''

###
# Create slides
###

lyt=prs.slide_layouts[0] # choosing a slide layout

# Simple slide
slide=prs.slides.add_slide(lyt) # adding a slide

title=slide.shapes.title # assigning a title
subtitle=slide.placeholders[1] # placeholder for subtitle

title.text="Hey,This is a Slide! How exciting!" # title
subtitle.text="Really?" # subtitle

# Blank slide + add image
blank_slide_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(blank_slide_layout)

image_path = "./data/demo.png"

# Position image
left=Inches(1)
top=Inches(0.5)

# Add image
img=slide2.shapes.add_picture(image_path,left,top)

prs.save("slide1.pptx") # saving file

###
# Load existing Slide deck
###

#prs2 = Presentation('existing-prs-file.pptx')
#prs2.save('new-file-name.pptx')